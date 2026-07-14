import subprocess
from pathlib import Path

import docx
import openpyxl
import pptx
import pymupdf
import pytest
from odf.opendocument import OpenDocumentSpreadsheet, OpenDocumentText
from odf.style import ParagraphProperties, Style, TableColumnProperties
from odf.table import Table, TableCell, TableColumn, TableRow
from odf.text import P
from openpyxl.utils import get_column_letter
from PIL import Image


@pytest.fixture
def make_image():
    """Factory fixture: writes a small test image to disk and returns its path."""

    def _make(
        path: Path, *, format: str, mode: str = "RGB", size: tuple[int, int] = (16, 12)
    ):
        image = Image.new(mode, size)
        for x in range(size[0]):
            for y in range(size[1]):
                if mode == "RGBA":
                    image.putpixel((x, y), ((x * 7) % 256, (y * 13) % 256, 128, 200))
                elif mode == "L":
                    image.putpixel((x, y), (x * 7 + y * 13) % 256)
                else:
                    image.putpixel((x, y), ((x * 7) % 256, (y * 13) % 256, 128))
        path.parent.mkdir(parents=True, exist_ok=True)
        image.save(path, format=format)
        return path

    return _make


@pytest.fixture
def make_pdf():
    """Factory fixture: writes a small multi-page test PDF to disk and returns its path."""

    def _make(
        path: Path,
        *,
        page_colors: tuple[tuple[float, float, float], ...] = ((1, 0, 0),),
        text: str | None = None,
        size: tuple[float, float] = (200, 300),
    ):
        path.parent.mkdir(parents=True, exist_ok=True)
        with pymupdf.open() as doc:
            for color in page_colors:
                page = doc.new_page(width=size[0], height=size[1])
                page.draw_rect(page.rect, color=color, fill=color)
                if text:
                    page.insert_text((20, 150), text, fontsize=24, color=(0, 0, 0))
            doc.save(path)
        return path

    return _make


@pytest.fixture
def make_docx():
    """Factory fixture: writes a small test .docx to disk and returns its path."""

    def _make(
        path: Path,
        *,
        paragraphs: tuple[str, ...] = ("Hello world",),
        page_break: bool = False,
    ):
        path.parent.mkdir(parents=True, exist_ok=True)
        document = docx.Document()
        for i, text in enumerate(paragraphs):
            if page_break and i > 0:
                document.add_page_break()
            document.add_paragraph(text)
        document.save(str(path))
        return path

    return _make


@pytest.fixture
def make_pptx():
    """Factory fixture: writes a small test .pptx to disk and returns its path."""

    def _make(path: Path, *, slide_texts: tuple[str, ...] = ("Hello world",)):
        path.parent.mkdir(parents=True, exist_ok=True)
        presentation = pptx.Presentation()
        layout = presentation.slide_layouts[1]  # title + content
        for text in slide_texts:
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = text
        presentation.save(str(path))
        return path

    return _make


@pytest.fixture
def make_odt():
    """Factory fixture: writes a small test .odt to disk and returns its path."""

    def _make(
        path: Path,
        *,
        paragraphs: tuple[str, ...] = ("Hello world",),
        page_break: bool = False,
    ):
        path.parent.mkdir(parents=True, exist_ok=True)
        document = OpenDocumentText()

        pagebreak_style = Style(name="PageBreak", family="paragraph")
        pagebreak_style.addElement(ParagraphProperties(breakbefore="page"))
        document.automaticstyles.addElement(pagebreak_style)

        for i, text in enumerate(paragraphs):
            paragraph = P(text=text)
            if page_break and i > 0:
                paragraph.setAttribute("stylename", pagebreak_style)
            document.text.addElement(paragraph)

        document.save(str(path))
        return path

    return _make


@pytest.fixture
def make_ods():
    """Factory fixture: writes a small test .ods to disk and returns its path."""

    def _make(
        path: Path,
        *,
        columns: int = 4,
        rows: int = 3,
        column_widths_mm: tuple[float, ...] | None = None,
    ):
        path.parent.mkdir(parents=True, exist_ok=True)
        document = OpenDocumentSpreadsheet()
        table = Table(name="Sheet1")

        for col in range(columns):
            table_column = TableColumn()
            if column_widths_mm is not None:
                column_style = Style(name=f"PetrificusCol{col}", family="table-column")
                column_style.addElement(
                    TableColumnProperties(columnwidth=f"{column_widths_mm[col]}mm")
                )
                document.automaticstyles.addElement(column_style)
                table_column.setAttribute("stylename", column_style)
            table.addElement(table_column)

        for row in range(1, rows + 1):
            table_row = TableRow()
            for col in range(1, columns + 1):
                cell = TableCell(valuetype="string")
                cell.addElement(P(text=f"r{row}c{col}"))
                table_row.addElement(cell)
            table.addElement(table_row)

        document.spreadsheet.addElement(table)
        document.save(str(path))
        return path

    return _make


@pytest.fixture
def make_text_file():
    """Factory fixture: writes a small test .txt file to disk and returns its path."""

    def _make(path: Path, *, content: str = "Hello world\n"):
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(content)
        return path

    return _make


@pytest.fixture
def make_csv():
    """Factory fixture: writes a small test .csv file to disk and returns its path."""

    def _make(
        path: Path,
        *,
        rows: tuple[tuple[str, ...], ...] = (("a", "b", "c"), ("1", "2", "3")),
    ):
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text("\n".join(",".join(row) for row in rows) + "\n")
        return path

    return _make


@pytest.fixture
def make_unsupported():
    """Factory fixture: writes a file with no registered CDR handler to disk.

    Content is arbitrary bytes rather than text, since dispatch is by
    sniffed content and plain text is a supported MIME type.
    """

    def _make(path: Path):
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(bytes(range(256)))
        return path

    return _make


@pytest.fixture
def make_xlsx():
    """Factory fixture: writes a small test .xlsx to disk and returns its path."""

    def _make(
        path: Path,
        *,
        columns: int = 4,
        rows: int = 3,
        column_widths: tuple[float, ...] | None = None,
    ):
        path.parent.mkdir(parents=True, exist_ok=True)
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        for row in range(1, rows + 1):
            for col in range(1, columns + 1):
                sheet.cell(row=row, column=col, value=f"r{row}c{col}")
        if column_widths is not None:
            for col, width in enumerate(column_widths, start=1):
                sheet.column_dimensions[get_column_letter(col)].width = width
        workbook.save(path)
        return path

    return _make


@pytest.fixture
def make_video(tmp_path):
    """Factory fixture: writes a small test video to disk using ffmpeg and returns its path.

    Video is a single solid color (so a specific pixel's value is a simple,
    deterministic thing to assert on) and audio is a pure sine tone.
    """

    def _make(
        path: Path,
        *,
        video_codec: str = "libx264",
        audio_codec: str | None = "aac",
        subtitles: tuple[str, ...] | None = None,
        subtitle_codec: str = "srt",
        color: str = "red",
        duration: float = 0.5,
        size: tuple[int, int] = (64, 48),
        rate: int = 5,
    ):
        path.parent.mkdir(parents=True, exist_ok=True)
        args = [
            "ffmpeg",
            "-y",
            "-f",
            "lavfi",
            "-i",
            f"color=c={color}:size={size[0]}x{size[1]}:duration={duration}:rate={rate}",
        ]
        if audio_codec:
            args += ["-f", "lavfi", "-i", f"sine=frequency=440:duration={duration}"]
        if subtitles:
            srt_path = tmp_path / f"{path.name}.srt"
            lines = []
            for i, text in enumerate(subtitles):
                start = f"00:00:{i:02d},000"
                end = f"00:00:{i + 1:02d},000"
                lines.append(f"{i + 1}\n{start} --> {end}\n{text}\n")
            srt_path.write_text("\n".join(lines))
            args += ["-i", str(srt_path)]

        args += ["-c:v", video_codec, "-pix_fmt", "yuv420p"]
        if audio_codec:
            args += ["-c:a", audio_codec]
        if subtitles:
            args += ["-c:s", subtitle_codec]
        args += [str(path)]

        subprocess.run(args, check=True, capture_output=True, timeout=30)
        return path

    return _make
